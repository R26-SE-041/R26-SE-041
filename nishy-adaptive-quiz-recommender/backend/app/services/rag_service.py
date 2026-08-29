"""
RAG Service — Document ingestion, chunking, embedding (sentence-transformers), and retrieval.
Uses ChromaDB as the vector store. Embeddings are computed locally (no API key needed).
"""
import os
import uuid
import logging
from pathlib import Path
from typing import List, Dict

import chromadb
from chromadb.config import Settings
from dotenv import load_dotenv

from app.services.llm_service import EmbeddingService

load_dotenv()
logger = logging.getLogger(__name__)

CHROMA_DIR    = os.getenv("CHROMA_PERSIST_DIR", "./db/chroma")
CHUNK_SIZE    = 800   # characters
CHUNK_OVERLAP = 150


class RagService:
    """
    Handles all document processing and vector search.
    Embeddings: sentence-transformers/all-MiniLM-L6-v2 (local, free).
    """

    def __init__(self):
        os.makedirs(CHROMA_DIR, exist_ok=True)
        self.client = chromadb.PersistentClient(
            path=CHROMA_DIR,
            settings=Settings(anonymized_telemetry=False)
        )
        self.embed = EmbeddingService()
        logger.info(f"RagService initialized | chroma_dir={CHROMA_DIR}")

    # ──────────────────────────────────────────────────────────
    # DOCUMENT EXTRACTION
    # ──────────────────────────────────────────────────────────

    def extract_pdf(self, file_path: str) -> List[Dict]:
        """Extract text chunks from PDF preserving headings and pages."""
        import fitz  # PyMuPDF
        doc = fitz.open(file_path)
        chunks = []
        current_heading = "Introduction"

        for page_num, page in enumerate(doc):
            blocks = page.get_text("dict")["blocks"]
            page_text_parts = []

            for block in blocks:
                if block["type"] == 0:
                    for line in block["lines"]:
                        for span in line["spans"]:
                            text = span["text"].strip()
                            if not text:
                                continue
                            if span["size"] > 14 and len(text) < 100:
                                current_heading = text
                            page_text_parts.append(text)

            page_text = " ".join(page_text_parts)
            if page_text.strip():
                for i, chunk in enumerate(self._chunk_text(page_text)):
                    chunks.append({
                        "chunk_id": f"{Path(file_path).stem}_p{page_num}_{i}",
                        "text": chunk,
                        "source": Path(file_path).name,
                        "page": page_num + 1,
                        "heading": current_heading,
                    })
        doc.close()
        logger.info(f"PDF extracted: {len(chunks)} chunks from {file_path}")
        return chunks

    def extract_docx(self, file_path: str) -> List[Dict]:
        """Extract text from DOCX file."""
        from docx import Document
        doc = Document(file_path)
        full_text = []
        current_heading = "Content"

        for para in doc.paragraphs:
            if para.style.name.startswith("Heading"):
                current_heading = para.text
            if para.text.strip():
                full_text.append(para.text)

        text = "\n".join(full_text)
        return [
            {
                "chunk_id": f"{Path(file_path).stem}_{i}",
                "text": chunk,
                "source": Path(file_path).name,
                "page": 0,
                "heading": current_heading,
            }
            for i, chunk in enumerate(self._chunk_text(text))
        ]

    def extract_pptx(self, file_path: str) -> List[Dict]:
        """Extract text from PPTX slides and speaker notes."""
        from pptx import Presentation
        prs = Presentation(file_path)
        chunks = []

        for slide_num, slide in enumerate(prs.slides):
            parts = []
            slide_title = f"Slide {slide_num + 1}"

            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    if shape.shape_type == 13:
                        slide_title = shape.text
                    parts.append(shape.text)

            if slide.has_notes_slide:
                notes = slide.notes_slide.notes_text_frame.text
                if notes.strip():
                    parts.append(f"[Notes]: {notes}")

            slide_text = " ".join(parts)
            if slide_text.strip():
                chunks.append({
                    "chunk_id": f"{Path(file_path).stem}_slide{slide_num + 1}",
                    "text": slide_text,
                    "source": Path(file_path).name,
                    "page": slide_num + 1,
                    "heading": slide_title,
                })
        return chunks

    def extract_txt(self, file_path: str) -> List[Dict]:
        """Extract text from plain text file."""
        with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
            text = f.read()
        return [
            {
                "chunk_id": f"{Path(file_path).stem}_{i}",
                "text": chunk,
                "source": Path(file_path).name,
                "page": 0,
                "heading": "Content",
            }
            for i, chunk in enumerate(self._chunk_text(text))
        ]

    # ──────────────────────────────────────────────────────────
    # CHUNKING
    # ──────────────────────────────────────────────────────────

    def _chunk_text(self, text: str) -> List[str]:
        """Split text into overlapping character chunks."""
        text = text.strip()
        if not text:
            return []
        chunks = []
        start = 0
        while start < len(text):
            chunk = text[start: start + CHUNK_SIZE]
            if chunk.strip():
                chunks.append(chunk)
            start += CHUNK_SIZE - CHUNK_OVERLAP
        return chunks

    # ──────────────────────────────────────────────────────────
    # CHROMADB OPERATIONS
    # ──────────────────────────────────────────────────────────

    def embed_and_store(self, collection_id: str, chunks: List[Dict]) -> None:
        """
        Batch-embed all chunks with sentence-transformers and store in ChromaDB.
        Much faster than one-by-one embedding.
        """
        if not chunks:
            logger.warning("embed_and_store called with empty chunks")
            return

        collection = self.client.get_or_create_collection(
            name=collection_id,
            metadata={"hnsw:space": "cosine"}
        )

        # Batch embed — sentence-transformers is fast locally
        BATCH = 64
        for i in range(0, len(chunks), BATCH):
            batch = chunks[i: i + BATCH]
            ids        = [c["chunk_id"] for c in batch]
            texts      = [c["text"] for c in batch]
            metadatas  = [
                {"source": c["source"], "page": c["page"], "heading": c["heading"]}
                for c in batch
            ]
            embeddings = self.embed.get_batch_embeddings(texts)

            collection.add(
                ids=ids,
                embeddings=embeddings,
                documents=texts,
                metadatas=metadatas,
            )
            logger.debug(f"Stored batch {i // BATCH + 1} | {len(batch)} chunks")

        logger.info(f"Stored {len(chunks)} chunks in collection '{collection_id}'")

    def retrieve(self, collection_id: str, query: str, k: int = 4) -> List[Dict]:
        """Retrieve top-k relevant chunks for a query using cosine similarity."""
        try:
            collection = self.client.get_collection(name=collection_id)
        except Exception:
            logger.warning(f"Collection '{collection_id}' not found")
            return []

        query_embedding = self.embed.get_query_embedding(query)
        n = min(k, collection.count())
        if n == 0:
            return []

        results = collection.query(
            query_embeddings=[query_embedding],
            n_results=n,
            include=["documents", "metadatas", "distances"],
        )

        chunks = []
        for i, doc in enumerate(results["documents"][0]):
            chunks.append({
                "chunk_id": results["ids"][0][i],
                "text": doc,
                "source": results["metadatas"][0][i].get("source", ""),
                "page": results["metadatas"][0][i].get("page", 0),
                "heading": results["metadatas"][0][i].get("heading", ""),
                "distance": results["distances"][0][i],
            })

        logger.debug(f"Retrieved {len(chunks)} chunks for: {query[:60]}")
        return chunks

    def get_source_chunks(self, collection_id: str, limit: int = 40) -> List[Dict]:
        """Read source chunks with metadata without inventing a semantic query."""
        try:
            collection = self.client.get_collection(name=collection_id)
        except Exception:
            logger.warning("Collection '%s' not found", collection_id)
            return []
        count = min(max(limit, 0), collection.count())
        if count == 0:
            return []
        results = collection.get(limit=count, include=["documents", "metadatas"])
        return [
            {
                "chunk_id": results["ids"][index],
                "text": document,
                "source": results["metadatas"][index].get("source", ""),
                "page": results["metadatas"][index].get("page", 0),
                "heading": results["metadatas"][index].get("heading", ""),
                "distance": 0.0,
            }
            for index, document in enumerate(results["documents"])
        ]

    def delete_collection(self, collection_id: str) -> None:
        """Delete a session's ChromaDB collection (cleanup)."""
        try:
            self.client.delete_collection(name=collection_id)
            logger.info(f"Deleted collection '{collection_id}'")
        except Exception as e:
            logger.warning(f"Could not delete collection '{collection_id}': {e}")

    def merge_collections(self, target_collection_id: str, source_collection_ids: List[str]) -> int:
        """Copy all chunks from multiple source collections into a target collection.

        Re-embeds the documents instead of fetching stored embeddings, which is
        more reliable across ChromaDB versions and avoids empty-embedding errors.
        """
        target = self.client.get_or_create_collection(
            name=target_collection_id,
            metadata={"hnsw:space": "cosine"}
        )
        total_copied = 0
        failures = []
        for src_id in source_collection_ids:
            try:
                src = self.client.get_collection(name=src_id)
                # Fetch documents and metadata only — re-embed locally to avoid
                # ChromaDB embedding-fetch failures on some versions.
                data = src.get(include=["documents", "metadatas"])
                if not data["ids"]:
                    logger.warning(f"Collection '{src_id}' exists but has no documents")
                    continue

                # Re-embed in batches and add to the target collection
                BATCH = 64
                ids = data["ids"]
                documents = data["documents"]
                metadatas = data["metadatas"]
                for i in range(0, len(ids), BATCH):
                    batch_ids = ids[i: i + BATCH]
                    batch_docs = documents[i: i + BATCH]
                    batch_meta = metadatas[i: i + BATCH]
                    # Skip IDs that are already present in the target collection
                    existing = set(target.get(ids=batch_ids)["ids"])
                    new_mask = [j for j, doc_id in enumerate(batch_ids) if doc_id not in existing]
                    if not new_mask:
                        total_copied += len(batch_ids)  # count as copied (already there)
                        continue
                    new_ids = [batch_ids[j] for j in new_mask]
                    new_docs = [batch_docs[j] for j in new_mask]
                    new_meta = [batch_meta[j] for j in new_mask]
                    embeddings = self.embed.get_batch_embeddings(new_docs)
                    target.add(
                        ids=new_ids,
                        embeddings=embeddings,
                        documents=new_docs,
                        metadatas=new_meta,
                    )
                    total_copied += len(new_ids)
                logger.info(f"Merged {len(ids)} chunks from '{src_id}' into '{target_collection_id}'")
            except Exception as e:
                logger.warning(f"Failed to copy from '{src_id}': {e}")
                failures.append(f"{src_id}: {e}")
        if total_copied == 0 and failures:
            raise RuntimeError(
                "Selected document collection(s) could not be read: " + "; ".join(failures)
            )
        logger.info(f"Merged {total_copied} chunks total into '{target_collection_id}'")
        return total_copied
