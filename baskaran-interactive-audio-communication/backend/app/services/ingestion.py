"""
Document ingestion service — Multilingual Hybrid RAG + Reranker.

Retrieval pipeline (hybrid_query_chunks):
  1. Dense Search  — ChromaDB cosine-similarity (multilingual embeddings)
  2. Sparse Search — BM25 keyword matching (rank-bm25)
  3. RRF Fusion    — Reciprocal Rank Fusion merges both ranked lists
  4. Reranker      — BAAI/bge-reranker-v2-m3 scores each candidate
  5. Return top-N  — reranked, highest-quality chunks only

Modal GPU mode (USE_MODAL_RETRIEVAL_MODELS=true):
  BGE-M3 embedding and reranking are offloaded to Modal GPU endpoints.
  When false, the existing local CPU path is used unchanged.

Why Hybrid:
  Tamil query "நாய் உணவு" → Dense may miss exact "dog food" keywords.
  BM25 catches keyword overlaps that embeddings miss.
  Reranker filters garbage — only truly relevant chunks reach the LLM.

Supports: PDF, PPTX, DOCX, XLSX, TXT, MD.
"""

import asyncio
import concurrent.futures
import contextlib
import hashlib
import io
import os
import pathlib
import re
import threading
import time
from typing import Generator

# A partial Hugging Face cache must fail promptly instead of leaving a web
# request waiting on a model download with no progress or response.
os.environ.setdefault("HF_HUB_ETAG_TIMEOUT", "10")
os.environ.setdefault("HF_HUB_DOWNLOAD_TIMEOUT", "60")
# BGE-M3's multi-gigabyte weights are fetched more reliably through the
# standard Hugging Face download path in this local Windows environment.
os.environ.setdefault("HF_HUB_DISABLE_XET", "1")

from sentence_transformers import SentenceTransformer, CrossEncoder

from app.db.chroma import get_or_create_collection
from app.core.config import get_settings
from app.core.logging import get_logger

logger = get_logger(__name__)

# ── Dedicated ML thread pool ──────────────────────────────────────────────────
# BGE-M3 encoding and CrossEncoder reranking are heavy CPU/memory operations.
# Using a dedicated executor instead of the default asyncio thread pool prevents
# a slow embedding from starving ALL other async tasks (including /health and
# simple list calls). max_workers=4 lets embedding + reranking run concurrently
# without blocking the server's general request handling.
_ML_EXECUTOR = concurrent.futures.ThreadPoolExecutor(
    max_workers=4, thread_name_prefix="ml_worker"
)


# ── Timing utility ────────────────────────────────────────────────────────────

@contextlib.contextmanager
def _timed(label: str):
    """
    Context manager that emits a [TIMING] log line on exit.

    Usage::

        with _timed("chroma_dense_search"):
            results = await ...
    """
    t0 = time.perf_counter()
    try:
        yield
    finally:
        elapsed = time.perf_counter() - t0
        logger.info("[TIMING] %s: %.3fs", label, elapsed)


class DocumentIngestionError(ValueError):
    """Raised when an uploaded document has no extractable text."""


class DocumentIngestionTimeoutError(DocumentIngestionError):
    """Raised when one ingestion stage exceeds the configured request limit."""


# ── Embedding model (Dense) ───────────────────────────────────────────────────
# Multilingual model — supports 50+ languages including Tamil & Sinhala.
EMBEDDING_MODEL = "BAAI/bge-m3"
EMBEDDING_DIMENSION = 1024
_embedder: SentenceTransformer | None = None
_embedder_lock = threading.Lock()


def _get_embedder() -> SentenceTransformer:
    """Load the one embedding model used for both documents and queries."""
    global _embedder
    if _embedder is None:
        with _embedder_lock:
            if _embedder is None:
                # Never download model weights from a request handler. The
                # image/local setup must provision BGE-M3 ahead of time; if it
                # is missing, return a diagnostic API error immediately.
                try:
                    _embedder = SentenceTransformer(EMBEDDING_MODEL, local_files_only=True)
                except Exception as exc:
                    logger.exception("BGE-M3 is unavailable in the local model cache")
                    raise DocumentIngestionError(
                        "BGE-M3 model files are not available locally, so this document "
                        "cannot be indexed yet. Complete the BAAI/bge-m3 model setup and retry."
                    ) from exc
                logger.info("Loaded multilingual embedder: %s", EMBEDDING_MODEL)
    return _embedder


def _encode_texts(texts: list[str]) -> list[list[float]]:
    """Produce normalized BGE-M3 embeddings in the collection's 1024-D space."""
    embeddings = _get_embedder().encode(
        texts, show_progress_bar=False, normalize_embeddings=True
    ).tolist()
    if any(len(vector) != EMBEDDING_DIMENSION for vector in embeddings):
        raise RuntimeError(
            f"{EMBEDDING_MODEL} returned an unexpected embedding dimension; "
            f"expected {EMBEDDING_DIMENSION}."
        )
    return embeddings


# ── Cross-Encoder Reranker ────────────────────────────────────────────────────
# BGE's multilingual reranker scores cross-language query/document pairs.
# It is lazy-loaded to avoid a cold-start cost for upload-only requests.
_reranker: CrossEncoder | None = None
_reranker_lock = threading.Lock()


def _get_reranker() -> CrossEncoder | None:
    """Lazy-load the cross-encoder reranker (singleton, thread-safe)."""
    global _reranker
    if _reranker is not None:
        return _reranker
    with _reranker_lock:
        if _reranker is None:
            try:
                _reranker = CrossEncoder(
                    "BAAI/bge-reranker-v2-m3",
                    max_length=512,
                    # Note: local_files_only=True breaks sharded safetensors weights
                    # (model.safetensors.index.json + shards). Let CrossEncoder find
                    # the cached files normally — weights load from disk, HF Hub only
                    # fetches a tiny metadata JSON if needed.
                )
                logger.info("Loaded multilingual reranker: BAAI/bge-reranker-v2-m3")
            except Exception as e:
                logger.warning("Reranker unavailable (%s) — skipping reranking step.", e)
    return _reranker


# ── BM25 Index (Sparse Search) ────────────────────────────────────────────────
# Per-user in-memory BM25 index.
# Rebuilt from ChromaDB documents on first query (or after new ingest).
# Thread-safe via lock.

class _BM25Store:
    """
    In-memory BM25 index for sparse keyword search.

    Stores all document chunks (text + metadata) per user_id.
    BM25 index is rebuilt lazily whenever the corpus changes.
    """

    def __init__(self):
        self._lock = threading.Lock()
        # { user_id: { "corpus": [str], "meta": [dict], "bm25": BM25Okapi | None, "dirty": bool } }
        self._store: dict[str, dict] = {}

    def _tokenize(self, text: str) -> list[str]:
        """Simple whitespace + punctuation tokenizer for BM25."""
        return re.findall(r"\b\w+\b", text.lower())

    def add_documents(self, user_id: str, texts: list[str], metadatas: list[dict]):
        """Add new document chunks to the BM25 corpus for a user."""
        with self._lock:
            if user_id not in self._store:
                self._store[user_id] = {"corpus": [], "meta": [], "bm25": None, "dirty": True}
            entry = self._store[user_id]
            entry["corpus"].extend(texts)
            entry["meta"].extend(metadatas)
            entry["dirty"] = True  # mark for rebuild on next query

    def _rebuild(self, user_id: str):
        """Rebuild BM25 index for a user (call with lock held)."""
        try:
            from rank_bm25 import BM25Okapi
            entry = self._store[user_id]
            tokenized = [self._tokenize(t) for t in entry["corpus"]]
            entry["bm25"] = BM25Okapi(tokenized)
            entry["dirty"] = False
            logger.debug("BM25 index rebuilt for user=%s (%d docs)", user_id, len(entry["corpus"]))
        except ImportError:
            logger.warning("rank-bm25 not installed — BM25 search disabled. Run: pip install rank-bm25")
        except Exception as e:
            logger.warning("BM25 index rebuild failed: %s", e)

    def search(self, user_id: str, query: str, top_k: int = 20) -> list[dict]:
        """Return top-K chunks by BM25 score for the given query."""
        with self._lock:
            if user_id not in self._store or not self._store[user_id]["corpus"]:
                return []
            entry = self._store[user_id]
            if entry["dirty"] or entry["bm25"] is None:
                self._rebuild(user_id)
            bm25 = entry["bm25"]
            if bm25 is None:
                return []
            tokens = self._tokenize(query)
            scores = bm25.get_scores(tokens)
            # Get top-K indices sorted by score descending
            top_indices = sorted(range(len(scores)), key=lambda i: scores[i], reverse=True)[:top_k]
            results = []
            for idx in top_indices:
                if scores[idx] > 0:   # skip zero-score (no keyword overlap)
                    results.append({
                        "text":     entry["corpus"][idx],
                        "metadata": entry["meta"][idx],
                        "bm25_score": float(scores[idx]),
                    })
            return results

    def load_from_chroma(self, user_id: str, documents: list[str], metadatas: list[dict]):
        """Populate BM25 index from ChromaDB dump (called at startup/first query)."""
        with self._lock:
            self._store[user_id] = {
                "corpus": list(documents),
                "meta":   list(metadatas),
                "bm25":   None,
                "dirty":  True,
            }


# Global BM25 store — shared across requests
_bm25_store = _BM25Store()


# ── RRF Fusion ────────────────────────────────────────────────────────────────

def _reciprocal_rank_fusion(
    dense_results: list[dict],
    sparse_results: list[dict],
    k: int = 60,
    dense_weight: float = 2.0,
    sparse_weight: float = 1.0,
) -> list[dict]:
    """
    Merge dense and sparse ranked lists using Reciprocal Rank Fusion.

    RRF score: score(d) = Σ 1 / (k + rank(d))  for each result list.
    k=60 is the standard constant from the original RRF paper (Cormack 2009).

    Returns merged list sorted by RRF score (descending), deduplicated by text.
    """
    scores: dict[str, float] = {}
    texts: dict[str, dict] = {}   # text → chunk dict (for dedup)

    for rank, chunk in enumerate(dense_results):
        key = chunk["text"][:100]   # use first 100 chars as dedup key
        scores[key] = scores.get(key, 0.0) + dense_weight / (k + rank + 1)
        texts[key] = chunk

    for rank, chunk in enumerate(sparse_results):
        key = chunk["text"][:100]
        scores[key] = scores.get(key, 0.0) + sparse_weight / (k + rank + 1)
        if key not in texts:
            texts[key] = chunk

    merged = sorted(scores.keys(), key=lambda k: scores[k], reverse=True)
    return [
        {**texts[k], "rrf_score": round(scores[k], 6)}
        for k in merged
    ]


# ── Text extractors ───────────────────────────────────────────────────────────

def _extract_pdf(file_bytes: bytes) -> list[dict]:
    """Extract text from PDF using PyMuPDF. Returns [{page, text}]."""
    import fitz  # PyMuPDF
    doc = fitz.open(stream=file_bytes, filetype="pdf")
    pages = []
    for i, page in enumerate(doc):
        text = page.get_text("text").strip()
        if text:
            pages.append({"page": i + 1, "text": text})
    return pages


def _extract_pptx(file_bytes: bytes) -> list[dict]:
    """Extract text from PowerPoint (.pptx). Each slide = one 'page'."""
    from pptx import Presentation
    prs = Presentation(io.BytesIO(file_bytes))
    pages = []
    for i, slide in enumerate(prs.slides):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = " ".join(run.text for run in para.runs).strip()
                    if line:
                        texts.append(line)
        combined = "\n".join(texts).strip()
        if combined:
            pages.append({"page": i + 1, "text": combined})
    return pages


def _extract_docx(file_bytes: bytes) -> list[dict]:
    """Extract text from Word (.docx). Returns as a single 'page'."""
    from docx import Document
    doc = Document(io.BytesIO(file_bytes))
    paragraphs = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
    if not paragraphs:
        return []
    pages = []
    for i in range(0, len(paragraphs), 20):
        chunk = "\n".join(paragraphs[i:i + 20])
        pages.append({"page": i // 20 + 1, "text": chunk})
    return pages


def _extract_xlsx(file_bytes: bytes) -> list[dict]:
    """Extract text from Excel (.xlsx). Each sheet = one 'page'."""
    import openpyxl
    wb = openpyxl.load_workbook(io.BytesIO(file_bytes), read_only=True, data_only=True)
    pages = []
    for sheet_idx, sheet in enumerate(wb.worksheets):
        rows = []
        for row in sheet.iter_rows(values_only=True):
            row_text = "\t".join(str(cell) for cell in row if cell is not None)
            if row_text.strip():
                rows.append(row_text)
        if rows:
            pages.append({"page": sheet_idx + 1, "text": "\n".join(rows)})
    return pages


def _extract_text_file(file_bytes: bytes) -> list[dict]:
    """Extract plain text / markdown. Returns as a single 'page'."""
    text = file_bytes.decode("utf-8", errors="replace").strip()
    if not text:
        return []
    return [{"page": 1, "text": text}]


# ── Dispatcher ────────────────────────────────────────────────────────────────

SUPPORTED_EXTENSIONS = {".pdf", ".pptx", ".docx", ".xlsx", ".txt", ".md"}

_EXTRACTORS = {
    ".pdf":  _extract_pdf,
    ".pptx": _extract_pptx,
    ".docx": _extract_docx,
    ".xlsx": _extract_xlsx,
    ".txt":  _extract_text_file,
    ".md":   _extract_text_file,
}


def _dispatch_extract(file_bytes: bytes, filename: str) -> list[dict]:
    """Choose the right extractor based on file extension."""
    ext = pathlib.Path(filename).suffix.lower()
    extractor = _EXTRACTORS.get(ext)
    if extractor is None:
        raise ValueError(f"Unsupported file type: '{ext}'. Supported: {sorted(SUPPORTED_EXTENSIONS)}")
    return extractor(file_bytes)


# ── Chunking ──────────────────────────────────────────────────────────────────

CHUNK_SIZE = 500    # characters
CHUNK_OVERLAP = 50


def _chunk_text(
    text: str,
    chunk_size: int = CHUNK_SIZE,
    overlap: int = CHUNK_OVERLAP,
) -> Generator[str, None, None]:
    """Sliding-window character-level chunking."""
    start = 0
    while start < len(text):
        end = min(start + chunk_size, len(text))
        yield text[start:end]
        start += chunk_size - overlap


# ── Public API ────────────────────────────────────────────────────────────────

async def ingest_document(
    file_bytes: bytes,
    document_id: str,
    filename: str,
    user_id: str,
) -> int:
    """
    Extract text from any supported document format, chunk it, embed it,
    and upsert into ChromaDB + BM25 index.

    Returns:
        Total number of chunks stored.
    """
    loop = asyncio.get_running_loop()
    timeout = get_settings().document_ingestion_timeout_seconds

    async def run_stage(stage: str, awaitable):
        logger.info("Document ingestion started: stage=%s file=%s", stage, filename)
        try:
            result = await asyncio.wait_for(awaitable, timeout=timeout)
        except asyncio.TimeoutError as exc:
            logger.error(
                "Document ingestion timed out: stage=%s file=%s timeout_seconds=%d",
                stage, filename, timeout,
            )
            raise DocumentIngestionTimeoutError(
                f"Indexing '{filename}' timed out during {stage}. Please retry."
            ) from exc
        logger.info("Document ingestion completed: stage=%s file=%s", stage, filename)
        return result

    pages = await run_stage(
        "text_extraction",
        loop.run_in_executor(None, _dispatch_extract, file_bytes, filename),
    )
    collection = await run_stage("chroma_collection_open", get_or_create_collection())

    chunks, ids, metadatas = [], [], []

    for page_info in pages:
        for chunk in _chunk_text(page_info["text"]):
            chunk = chunk.strip()
            if not chunk:
                continue
            chunk_id = hashlib.md5(f"{document_id}:{chunk}".encode()).hexdigest()
            chunks.append(chunk)
            ids.append(chunk_id)
            metadatas.append({
                "document_id": document_id,
                "filename": filename,
                "user_id": user_id,
                "page": page_info["page"],
                "embedding_model": EMBEDDING_MODEL,
            })

    if not chunks:
        logger.warning("No extractable text found in %s", filename)
        raise DocumentIngestionError(
            "No readable text was found in this document. "
            "If this is a scanned PDF, upload a text-searchable PDF or a DOCX/TXT version."
        )

    # Dense embeddings → ChromaDB
    # When USE_MODAL_RETRIEVAL_MODELS=true, route embedding to Modal GPU for
    # consistency with query-time embeddings (same model, same normalization).
    settings = get_settings()
    if settings.use_modal_retrieval_models:
        from app.services.modal_client import call_bge_embed
        embeddings = await run_stage(
            "bge_m3_embedding_modal",
            call_bge_embed(chunks),
        )
    else:
        embeddings = await run_stage(
            "bge_m3_embedding",
            loop.run_in_executor(_ML_EXECUTOR, lambda: _encode_texts(chunks)),
        )

    await run_stage(
        "chromadb_upsert",
        collection.upsert(
            ids=ids,
            documents=chunks,
            embeddings=embeddings,
            metadatas=metadatas,
        ),
    )

    # Sparse index → BM25 (also add to in-memory store)
    _bm25_store.add_documents(user_id, chunks, metadatas)

    logger.info("Ingested %d chunks from %s (ChromaDB + BM25)", len(chunks), filename)
    return len(chunks)


async def _ensure_bm25_loaded(user_id: str) -> None:
    """
    Populate the BM25 index from ChromaDB if it hasn't been loaded yet.
    Called lazily on first query after a server restart (BM25 is in-memory only).
    """
    # If BM25 store already has data for this user, skip
    if _bm25_store._store.get(user_id, {}).get("corpus"):
        return

    try:
        collection = await get_or_create_collection()
        loop = asyncio.get_event_loop()

        # Fetch ALL chunks for this user from ChromaDB to rebuild BM25 index
        count = await collection.count() if hasattr(collection, "count") else None
        if not count:
            return

        results = await collection.get(
            where={"user_id": user_id},
            include=["documents", "metadatas"],
        )

        docs = results.get("documents") or []
        metas = results.get("metadatas") or []

        if docs:
            _bm25_store.load_from_chroma(user_id, docs, metas)
            logger.info("BM25 index loaded from ChromaDB: %d chunks for user=%s", len(docs), user_id)
    except Exception as e:
        logger.warning("BM25 index warm-up failed (non-fatal): %s", e)


async def query_chunks(
    query: str,
    user_id: str,
    n_results: int = 5,
) -> list[dict]:
    """
    Legacy dense-only retrieval. Kept for backwards compatibility.
    Use hybrid_query_chunks() for better accuracy.
    """
    try:
        collection = await get_or_create_collection()
        loop = asyncio.get_event_loop()

        try:
            existing = await collection.count() if hasattr(collection, "count") else None
        except Exception:
            existing = None

        if existing is not None and existing == 0:
            return []

        safe_n = min(n_results, existing) if existing is not None else n_results
        if safe_n < 1:
            return []

        query_embedding = await asyncio.wait_for(
            loop.run_in_executor(_ML_EXECUTOR, lambda: _encode_texts([query])[0]),
            timeout=90.0,  # 90 s covers BGE-M3 first-load on CPU; 2-5 s after warm-up
        )

        results = await collection.query(
            query_embeddings=[query_embedding],
            n_results=safe_n,
            where={"user_id": user_id},
            include=["documents", "metadatas", "distances"],
        )

        chunks = []
        for doc, meta, dist in zip(
            (results.get("documents") or [[]])[0],
            (results.get("metadatas") or [[]])[0],
            (results.get("distances") or [[]])[0],
        ):
            chunks.append({
                "text":     doc,
                "metadata": meta,
                "score":    round(max(0.0, min(1.0, 1 - dist)), 4),
            })
        return chunks

    except Exception as e:
        logger.warning("query_chunks failed (returning empty): %s", e)
        return []


async def hybrid_query_chunks(
    query: str,
    user_id: str,
    n_results: int = 5,
    dense_candidates: int = 10,
    sparse_candidates: int = 10,
) -> list[dict]:
    """
    Multilingual Hybrid RAG + Reranker retrieval pipeline.

    Steps:
      1. Dense Search  — ChromaDB embedding similarity (top-10)
      2. Sparse Search — BM25 keyword matching (top-10)
      3. RRF Fusion    — Merge ranked lists (8-10 unique candidates)
      4. Reranker      — BGE cross-encoder scores each candidate vs query
                         (Modal GPU when USE_MODAL_RETRIEVAL_MODELS=true,
                          local CPU otherwise)
      5. Return top-N  — highest-quality, truly relevant chunks

    Args:
        query:             The search query (any language).
        user_id:           Filter results to this user's documents.
        n_results:         Final number of chunks to return (after reranking).
        dense_candidates:  How many chunks to fetch from ChromaDB.
        sparse_candidates: How many chunks to fetch from BM25.

    Returns:
        List of chunk dicts with keys: text, metadata, score, retrieval_method.
    """
    import time as _time
    loop = asyncio.get_event_loop()
    settings = get_settings()
    use_modal = settings.use_modal_retrieval_models

    # ── Step 1: Dense Search ──────────────────────────────────────────────────
    dense_results: list[dict] = []
    try:
        collection = await get_or_create_collection()

        try:
            existing = await collection.count() if hasattr(collection, "count") else None
        except Exception:
            existing = None

        if existing and existing > 0:
            safe_n = min(dense_candidates, existing)

            # ── Embedding: Modal GPU or local CPU ────────────────────────────
            if use_modal:
                from app.services.modal_client import call_bge_embed
                with _timed("modal_bge_embedding"):
                    embed_vecs = await call_bge_embed([query])
                query_embedding = embed_vecs[0]
            else:
                with _timed("local_bge_embedding"):
                    query_embedding = await asyncio.wait_for(
                        loop.run_in_executor(
                            _ML_EXECUTOR, lambda: _encode_texts([query])[0]
                        ),
                        timeout=90.0,  # 90 s covers BGE-M3 first-load on CPU
                    )

            with _timed("chroma_dense_search"):
                results = await collection.query(
                    query_embeddings=[query_embedding],
                    n_results=safe_n,
                    where={"user_id": user_id},
                    include=["documents", "metadatas", "distances"],
                )
            for doc, meta, dist in zip(
                (results.get("documents") or [[]])[0],
                (results.get("metadatas") or [[]])[0],
                (results.get("distances") or [[]])[0],
            ):
                dense_results.append({
                    "text":     doc,
                    "metadata": meta,
                    "score":    round(max(0.0, min(1.0, 1 - dist)), 4),
                    "retrieval_method": "dense",
                })
            logger.debug("Dense search: %d results", len(dense_results))
    except Exception as e:
        logger.warning("Dense search failed: %s", e)

    # ── Step 2: Sparse Search (BM25) ─────────────────────────────────────────
    sparse_results: list[dict] = []
    try:
        # Warm up BM25 index from ChromaDB if server just restarted
        await _ensure_bm25_loaded(user_id)
        with _timed("bm25_search"):
            raw_sparse = await loop.run_in_executor(
                _ML_EXECUTOR,
                lambda: _bm25_store.search(user_id, query, top_k=sparse_candidates),
            )
        for chunk in raw_sparse:
            sparse_results.append({
                "text":     chunk["text"],
                "metadata": chunk["metadata"],
                "score":    round(chunk["bm25_score"], 4),
                "retrieval_method": "bm25",
            })
        logger.debug("BM25 search: %d results", len(sparse_results))
    except Exception as e:
        logger.warning("BM25 search failed (non-fatal, dense-only fallback): %s", e)

    # ── Step 3: RRF Fusion ────────────────────────────────────────────────────
    with _timed("hybrid_merge"):
        if dense_results and sparse_results:
            fused = _reciprocal_rank_fusion(dense_results, sparse_results)
            logger.debug("RRF fusion: %d unique candidates", len(fused))
        elif dense_results:
            fused = dense_results   # BM25 unavailable — dense only
            logger.debug("Hybrid: BM25 unavailable, using dense-only results")
        else:
            fused = sparse_results  # ChromaDB unavailable — sparse only

    if not fused:
        logger.info("hybrid_query_chunks: no candidates found for user=%s", user_id)
        return []

    # Send up to 10 candidates to the reranker (dense_candidates + sparse_candidates
    # produces 8-10 unique chunks after dedup via RRF).
    candidates = fused[:10]

    # ── Step 4: Reranking — Modal GPU or local CPU ───────────────────────────
    if use_modal:
        # Modal GPU reranker — returns candidates sorted, already top_k sliced
        from app.services.modal_client import call_bge_rerank
        with _timed("modal_reranker"):
            candidates = await call_bge_rerank(query, candidates, top_k=n_results)
        # modal reranker already sliced to top_k; re-normalise score field
        for chunk in candidates:
            chunk["score"] = round(
                chunk.get("reranker_score", chunk.get("rrf_score", chunk.get("score", 0.0))),
                4,
            )
        return candidates
    else:
        # Local CPU CrossEncoder (original path — unchanged)
        reranker = _get_reranker()
        if reranker and len(candidates) > 1:
            try:
                pairs = [(query, c["text"]) for c in candidates]
                with _timed("local_reranker"):
                    scores = await asyncio.wait_for(
                        loop.run_in_executor(_ML_EXECUTOR, lambda: reranker.predict(pairs)),
                        timeout=20.0,  # reranker timeout; skip if slow
                    )
                for chunk, score in zip(candidates, scores):
                    logit = float(score)
                    chunk["reranker_score"] = 1.0 / (1.0 + pow(2.718281828, -logit))
                candidates.sort(key=lambda c: c.get("reranker_score", 0.0), reverse=True)
                logger.debug("Reranker scored %d candidates", len(candidates))
            except Exception as e:
                logger.warning("Reranker failed (non-fatal, using RRF order): %s", e)

        # ── Step 5: Return top-N ──────────────────────────────────────────────
        top = candidates[:n_results]
        for chunk in top:
            chunk["score"] = round(
                chunk.get("reranker_score", chunk.get("rrf_score", chunk.get("score", 0.0))),
                4,
            )
        return top


# ── Backwards-compat alias ────────────────────────────────────────────────────
async def _document_is_indexed(collection, document_id: str) -> bool:
    """Check the BGE-M3 collection before reindexing an original upload."""
    result = await collection.get(where={"document_id": document_id}, include=["metadatas"])
    return bool(result.get("ids"))


async def reindex_stored_documents() -> None:
    """Populate the versioned BGE-M3 collection from preserved original uploads.

    The legacy Chroma collection is deliberately left untouched: it is neither
    queried nor mixed with the new vectors. Reindexing is idempotent, so a
    restart resumes any interrupted migration without deleting PDFs.
    """
    from app.services import local_document_store

    collection = await get_or_create_collection()
    migrated = 0
    for record in local_document_store.list_all_documents():
        try:
            if await _document_is_indexed(collection, record["document_id"]):
                continue
            await ingest_document(
                local_document_store.read_document(record),
                record["document_id"],
                record["filename"],
                record["user_id"],
            )
            migrated += 1
        except FileNotFoundError:
            logger.warning("Cannot reindex missing local document %s", record["document_id"])
        except Exception as exc:
            logger.warning("Could not reindex local document %s: %s", record["document_id"], exc)

    # Production uploads are retained in Supabase Storage. This is best-effort
    # so an unavailable Supabase instance never prevents application startup.
    try:
        from app.db.supabase import get_supabase
        from app.services.storage import download_document

        client = await get_supabase()
        result = await client.table("documents").select(
            "id,user_id,filename,storage_path"
        ).execute()
        for record in result.data or []:
            if not record.get("storage_path") or record["storage_path"].startswith("local/"):
                continue
            if await _document_is_indexed(collection, record["id"]):
                continue
            try:
                content = await download_document(record["storage_path"])
                await ingest_document(content, record["id"], record["filename"], record["user_id"])
                migrated += 1
            except Exception as exc:
                logger.warning("Could not reindex Supabase document %s: %s", record["id"], exc)
    except Exception as exc:
        logger.info("Supabase BGE-M3 reindex skipped: %s", exc)

    if migrated:
        logger.info("BGE-M3 migration indexed %d stored document(s)", migrated)


ingest_pdf = ingest_document
